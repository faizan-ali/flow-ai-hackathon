from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Project Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "3D Visualization for Emergency Response"
subtitle.text = "Leveraging Text-to-3D Technology to Enhance Emergency Response"

# Slide 2: Problem & Solution
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Problem & Solution"
body.text = ("Problem:\n"
              "- Current emergency responses rely heavily on verbal descriptions and 2D maps\n"
              "- Communication barriers due to language and descriptive limitations\n\n"
              "Solution:\n"
              "- Text-to-3D visualization using GPT-4 and shap-e diffusion mode\n"
              "- Facilitating non-verbal communication through interactive 3D scenes")

# Slide 3: Innovation
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Innovation"
body.text = ("Unique Approach:\n"
              "- Novel text-to-3D visualization technology\n"
              "- Bridging communication gaps through intuitive communication methods\n\n"
              "Creative Use of AI:\n"
              "- Leveraging GPT-4 for intelligent text interpretation\n"
              "- Real-time interactive scenarios for dynamic response planning")

# Slide 4: Technical Excellence
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Technical Excellence"
body.text = ("Quality AI Models:\n"
              "- Spatial awareness capabilities of GPT-4\n"
              "- Domain-specific training for accurate 3D representations\n\n"
              "Relevant Technologies:\n"
              "- Integration with GIS and other databases for real-world data\n"
              "- Dynamic visualizations advancing beyond static 2D representations")

# Slide 5: Applications & Next Steps
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Applications & Next Steps"
body.text = ("Applications:\n"
              "- Emergency response and community engagement\n"
              "- Training simulations and post-emergency analysis\n\n"
              "Next Steps:\n"
              "- Further integration with real-world databases\n"
              "- Enhancing user experience with more interactive features")

# Save the presentation
presentation_path = '3D_Visualization_Emergency_Response_Presentation.pptx'
prs.save(presentation_path)

presentation_path
